"""Conversion smoke tests for automated conversion engine bumps.

This is the gate the automated engine-bump workflow relies on: if a new markitdown release
breaks a real conversion path, this must fail in CI *before* a release is built. It therefore
covers the binary office/PDF formats that are the engine's actual value — not just the trivial
text formats.

Fixtures are generated in-process rather than committed. That keeps the suite hermetic (no
generation libraries needed at run time beyond what the engine extras already pull in) and
avoids committing binary blobs into the repository.
"""

from __future__ import annotations

import io
import zipfile

from omnivert.conversion import service
from omnivert.schemas import ConvertOptions


def _check(name: str, result) -> None:
    if not result.ok or not result.markdown:
        raise AssertionError(f"{name} conversion failed: {result.error}")


def assert_text_ok(name: str, content: str, extension: str) -> None:
    _check(name, service.convert_text(content, extension, "utf-8", ConvertOptions()))


def assert_bytes_ok(name: str, data: bytes, filename: str) -> None:
    _check(name, service.convert_bytes(data, filename, ConvertOptions()))


# --- fixture builders ----------------------------------------------------------------

_SAMPLE = "Hello from Omnivert smoke test"


def _xlsx_bytes() -> bytes:
    """A one-cell workbook (openpyxl ships with the engine's xlsx extra)."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["name", "value"])
    ws.append(["alpha", 1])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    """A single title slide (python-pptx ships with the engine's pptx extra)."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = _SAMPLE
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _docx_bytes() -> bytes:
    """A minimal valid OOXML word document, built from stdlib zipfile so no writer lib is
    required (the engine reads .docx via mammoth, which has no document writer)."""
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{_SAMPLE}</w:t></w:r></w:p></w:body>"
        "</w:document>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
    return buf.getvalue()


def _pdf_bytes() -> bytes:
    """A minimal one-page PDF with a single text run, assembled by hand (no PDF writer lib
    ships with the engine's pdf extra — pdfminer only reads). Byte offsets feed a valid xref
    table so pdfminer can extract the text."""
    stream = b"BT /F1 24 Tf 72 720 Td (" + _SAMPLE.encode("ascii") + b") Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + body + b"\nendobj\n"
    xref_pos = len(out)
    size = len(objects) + 1
    out += b"xref\n0 " + str(size).encode() + b"\n"
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode("ascii")
    out += b"trailer\n<< /Size " + str(size).encode() + b" /Root 1 0 R >>\n"
    out += b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    return bytes(out)


def main() -> None:
    assert_text_ok("plain text", "hello world", ".txt")
    assert_text_ok("csv", "name,value\nalpha,1\n", ".csv")
    assert_text_ok("html", "<html><body><h1>Hello</h1><p>World</p></body></html>", ".html")
    assert_bytes_ok("xlsx", _xlsx_bytes(), "sample.xlsx")
    assert_bytes_ok("pptx", _pptx_bytes(), "sample.pptx")
    assert_bytes_ok("docx", _docx_bytes(), "sample.docx")
    assert_bytes_ok("pdf", _pdf_bytes(), "sample.pdf")
    print("engine smoke: all formats OK")


if __name__ == "__main__":
    main()
